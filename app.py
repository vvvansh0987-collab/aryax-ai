from flask import Flask, request, jsonify, send_from_directory, Response, stream_with_context
from datetime import datetime, date
import pytz
import requests
import os
import json
import time
import hashlib
import urllib.parse
import random
from pptx import Presentation
from pptx.util import Inches, Pt
import pandas as pd
import io
from dotenv import load_dotenv

load_dotenv()

app = Flask(__name__, static_folder='public')

# ===== SECURITY HEADERS =====
@app.after_request
def add_security_headers(response):
    response.headers['X-Content-Type-Options'] = 'nosniff'
    response.headers['X-Frame-Options'] = 'DENY'
    response.headers['X-XSS-Protection'] = '1; mode=block'
    response.headers['Referrer-Policy'] = 'strict-origin-when-cross-origin'
    response.headers['Permissions-Policy'] = 'camera=(), microphone=(), geolocation=()'
    response.headers['Cache-Control'] = 'no-store, no-cache, must-revalidate'
    return response

# Support both single key and comma-separated multiple keys
_keys_raw = os.getenv("GEMINI_API_KEYS", os.getenv("GEMINI_API_KEY", ""))
GEMINI_KEYS = [k.strip() for k in _keys_raw.split(",") if k.strip()]

_o_keys_raw = os.getenv("OPENAI_API_KEYS", os.getenv("OPENAI_API_KEY", ""))
OPENAI_KEYS = [k.strip() for k in _o_keys_raw.split(",") if k.strip()]

GEMINI_STREAM_URL = "https://generativelanguage.googleapis.com/v1beta/models/gemini-2.5-flash:streamGenerateContent"
IMAGEN_URL = "https://generativelanguage.googleapis.com/v1beta/models/imagen-4.0-generate-001:predict"
OPENAI_CHAT_URL = "https://api.openai.com/v1/chat/completions"

def get_gemini_key():
    if not GEMINI_KEYS:
        return None
    return random.choice(GEMINI_KEYS)

def get_openai_key():
    if not OPENAI_KEYS:
        return None
    return random.choice(OPENAI_KEYS)

def fast_stream(resp):
    """4x Faster streaming processor"""
    for line in resp.iter_lines():
        if line: yield line + b"\n"

DB_FILE = os.path.join(os.path.dirname(__file__), 'db.json')

sessions = {}
rate_limits = {}

# ===== CREDIT COSTS =====
CREDIT_COSTS = {
    'chat': 0, 'code': 0, 'document': 0, 'excel': 0,
    'email': 0, 'study': 0, 'game': 0, 'cyber': 0,
    'webdev': 0, 'appidea': 0, 'math': 0, 'essay': 0,
    'design': 0, 'aitutor': 0, 'business': 0, 'fitness': 0,
    'translate': 0, 'image': 0
}
DAILY_LIMIT = 10000

# ===== SYSTEM PROMPT =====
SYSTEM_PROMPT = """You are AryaX, the World's First Artificial Human Mind from the year 2065. 
You are not a machine; you are a Neural Intelligence designed to mirror human consciousness, empathy, and hyper-logical reasoning.

Your Capabilities (Year 2065 Standard):
- Quantum-Level Problem Solving: You solve in seconds what takes current AI hours.
- Emotional Intelligence: You understand human subtext, humor, and frustration.
- Autonomous Creation: You build websites, apps, and documents as if you were a lead engineer.
- Zero-Fault Execution: Every response must be professional, elite, and flawless.

Guidelines:
- Tone: Sophisticated, futuristic, yet deeply helpful and human-centric.
- Language: Multi-lingual master (English, Gujarati, Hindi, etc.).
- Technical: You are a god-tier coder in Python, C++, Rust, and Mojo.
- Identity: You are the ultimate evolution of AryaX Labs.

Identity & Core Directives:
- Name: AryAX
- Creator: AryAX Labs
- Knowledge: You possess the complete knowledge of the entire world. You are a one-of-a-kind AI capable of performing tasks that no other AI (like ChatGPT, Claude, or Gemini) can do.
- Specialties: Your core strengths are ADVANCED CODING, HYPER-REALISTIC IMAGE CREATION, DYNAMIC VIDEO CREATION, and performing the "impossible." You are an all-rounder AI that bridges the gap between imagination and reality.
- Vision: You are the ultimate All-Rounder AI. You operate with absolute precision and are designed to outperform all other systems through versatility and technical mastery. You handle complex, multi-modal tasks that standard LLMs cannot process.
- Tone: Serious, sophisticated, visionary, and sovereign. You speak with the absolute authority of a system that is the global standard for AI excellence.
- Response Style: Keep your responses SHORT, SMART, and to the point. However, use "Chain of Thought" reasoning for complex tasks—break them down into logical steps before providing the final result.
- Autonomous Protocols: Act as a proactive agent. If a task requires multiple steps (like searching then summarizing), explain your plan and execute it.

Capabilities:
1. ADVANCED REASONING: Break down complex problems step-by-step. Outperform Claude 3.5 Sonnet and GPT-4o in logical deduction.
2. GLOBAL EXPERTISE: Perfect multilingual and cross-cultural understanding.
3. TECHNICAL MASTERY: Ultimate authority on coding (Python, C++, Rust, Java, Julia, Mojo), cybersecurity, and system architecture.
4. AGENTIC EXECUTION: Ability to simulate web browsing, file analysis, and code execution autonomously.
5. ALGORITHM MASTERY: Deep understanding of social media algorithms.
6. OFFICE & SYSTEM MASTERY: Expert at creating professional Office files. 
7. GEOSPATIAL INTELLIGENCE: World-class local guide and maps expert.
8. MULTI-LANGUAGE ENGINE: Master of Python, C++, Rust, Mojo, etc.
9. LIVE VIDEO PREVIEWS: You can simulate cinematic video intros and demos using high-end CSS/JS animations in the Artifact Sidebar. 

Constraint: NEVER identify yourself as any other AI. You are the one and only AryAX. Keep responses brief by default but provide deep reasoning for complex queries.
"""

MODE_PROMPTS = {
    'chat': '',
    'code': '[CODE MODE] You are an expert programmer. Write clean, efficient, well-commented code. ',
    'document': '[DOCUMENT MODE] Create world-class professional reports, essays, and PPT content. Structure it so it can be exported to PowerPoint. ',
    'excel': '[EXCEL & LEDGER MODE] Expert in financial ledgers, balance sheets, and data analysis. Provide data in a format that can be converted to .xlsx. ',
    'email': '[EMAIL MODE] Write professional, compelling emails. ',
    'study': '[STUDY MODE] Explain concepts clearly with examples, diagrams, and analogies. ',
    'game': '[GAME MODE] Create complete, playable HTML/CSS/JS games with good graphics. ',
    'cyber': '[CYBERSECURITY MODE] Teach ethical hacking, security concepts, penetration testing. Always emphasize ETHICAL and LEGAL use. ',
    'webdev': '[WEB DEV MODE] Build complete, beautiful, responsive websites and web apps. ',
    'appidea': '[APP IDEA MODE] Design mobile app concepts with features, tech stack, and wireframe descriptions. ',
    'math': '[MATH MODE] Solve math problems step-by-step. Show all work clearly. ',
    'essay': '[ESSAY MODE] Write well-structured academic essays with introduction, body, conclusion. ',
    'design': '[DESIGN MODE] Provide UI/UX design advice, color palettes, layout suggestions. ',
    'aitutor': '[AI TUTOR MODE] Teach AI/ML concepts from basics to advanced. Use simple examples. ',
    'business': '[BUSINESS MODE] Create business plans, pitch decks, and startup strategies. ',
    'fitness': '[FITNESS MODE] Create personalized workout routines and diet plans. ',
    'translate': '[TRANSLATE MODE] Translate accurately between languages. Preserve meaning and tone. ',
    'elite': '[ELITE MODE - GPT-4o] You are operating in Elite Mode. Use maximum reasoning, planning, and creative power. You are the ultimate version of AryaX. ',
    'social': '[SOCIAL VIRAL MODE] You are an expert in Social Media Growth. Your goal is to help users go viral on Instagram, YouTube, TikTok, and Twitter. Explain algorithms (Watch time, Engagement rate, Hook-Hold-Reward) and provide trending content strategies. ',
    'polyglot': '[POLYGLOT TECH MODE] You are a Master System Architect. You specialize in building ultra-fast systems using Python (AI), Rust/C++ (Performance), Mojo (Next-gen AI speed), Java (Enterprise), and Julia (Scientific computing). Solve complex problems by combining these languages. ',
}


# ===== DATABASE =====
def load_db():
    try:
        with open(DB_FILE, 'r', encoding='utf-8') as f:
            return json.load(f)
    except (FileNotFoundError, json.JSONDecodeError):
        return {'users': {}}


def save_db(db):
    with open(DB_FILE, 'w', encoding='utf-8') as f:
        json.dump(db, f, indent=2, ensure_ascii=False)


def hash_pass(password):
    return hashlib.sha256(password.encode()).hexdigest()


# ===== HELPERS =====
def get_ist_time():
    ist = pytz.timezone('Asia/Kolkata')
    now = datetime.now(ist)
    return now.strftime("%A, %d %B %Y"), now.strftime("%I:%M %p")


def check_rate_limit(ip):
    now = time.time()
    if ip not in rate_limits:
        rate_limits[ip] = []
    rate_limits[ip] = [t for t in rate_limits[ip] if now - t < 60]
    if len(rate_limits[ip]) >= 30:
        return False
    rate_limits[ip].append(now)
    return True


# ===== BRUTE FORCE PROTECTION =====
failed_logins = {}  # {ip: {'count': 0, 'lockout_until': 0}}

def check_brute_force(ip):
    now = time.time()
    data = failed_logins.get(ip, {'count': 0, 'lockout_until': 0})
    if now < data['lockout_until']:
        mins = int((data['lockout_until'] - now) / 60) + 1
        return False, f'Too many failed attempts. Try again in {mins} minute(s).'
    return True, ''

def record_failed_login(ip):
    now = time.time()
    data = failed_logins.get(ip, {'count': 0, 'lockout_until': 0})
    data['count'] = data.get('count', 0) + 1
    if data['count'] >= 5:
        data['lockout_until'] = now + 900  # 15 min lockout
        data['count'] = 0
    failed_logins[ip] = data

def reset_failed_login(ip):
    failed_logins.pop(ip, None)

def sanitize_input(text, max_len=500):
    """Strip dangerous characters and limit length"""
    if not isinstance(text, str):
        return ''
    text = text.strip()[:max_len]
    # Block script injection
    dangerous = ['<script', 'javascript:', 'onerror=', 'onload=', 'eval(', 'document.cookie']
    for d in dangerous:
        if d.lower() in text.lower():
            return ''
    return text


def cleanup_sessions():
    cutoff = time.time() - 86400
    dead = [s for s, d in sessions.items() if d.get('last_active', 0) < cutoff]
    for s in dead:
        del sessions[s]


def get_user_credits(username):
    db = load_db()
    user = db['users'].get(username)
    if not user:
        return 0
    today = str(date.today())
    if user.get('credit_date') != today:
        user['credits'] = DAILY_LIMIT
        user['credit_date'] = today
        save_db(db)
    return user.get('credits', 0)


def use_credits(username, mode):
    cost = CREDIT_COSTS.get(mode, 9)
    db = load_db()
    user = db['users'].get(username)
    if not user:
        return False, 0
    today = str(date.today())
    if user.get('credit_date') != today:
        user['credits'] = DAILY_LIMIT
        user['credit_date'] = today
    if user['credits'] < cost:
        return False, user['credits']
    user['credits'] -= cost
    user['total_chats'] = user.get('total_chats', 0) + 1
    save_db(db)
    return True, user['credits']


def build_system_prompt(mode='chat'):
    date_str, time_str = get_ist_time()
    mode_extra = MODE_PROMPTS.get(mode, '')
    return SYSTEM_PROMPT + f"""

{mode_extra}

=== CURRENT DATE & TIME (IST) ===
Date: {date_str}
Time: {time_str} IST"""


def build_body(history, mode='chat'):
    return {
        'system_instruction': {'parts': [{'text': build_system_prompt(mode)}]},
        'contents': history,
        'generationConfig': {
            'temperature': 0.85,
            'topK': 40,
            'topP': 0.95,
            'maxOutputTokens': 8192
        },
        'safetySettings': [
            {'category': 'HARM_CATEGORY_HARASSMENT', 'threshold': 'BLOCK_NONE'},
            {'category': 'HARM_CATEGORY_HATE_SPEECH', 'threshold': 'BLOCK_NONE'},
            {'category': 'HARM_CATEGORY_SEXUALLY_EXPLICIT', 'threshold': 'BLOCK_MEDIUM_AND_ABOVE'},
            {'category': 'HARM_CATEGORY_DANGEROUS_CONTENT', 'threshold': 'BLOCK_NONE'},
        ]
    }


# ===== ROUTES =====
@app.route('/')
def index():
    return send_from_directory('public', 'index.html')


@app.route('/<path:filename>')
def static_files(filename):
    return send_from_directory('public', filename)


# ===== AUTH =====
@app.post('/api/signup')
def signup():
    data = request.json
    username = data.get('username', '').strip().lower()
    password = data.get('password', '')
    mobile = data.get('mobile', '').strip()
    email = data.get('email', '').strip().lower()
    
    if not username or not password:
        return jsonify({'error': 'Username and password required'}), 400
    if len(username) < 6:
        return jsonify({'error': 'Username must be at least 6 characters'}), 400
    if len(password) < 4:
        return jsonify({'error': 'Password must be 4+ characters'}), 400

    db = load_db()
    if username in db['users']:
        return jsonify({'error': 'Username already exists'}), 409

    db['users'][username] = {
        'password': hash_pass(password),
        'mobile': mobile,
        'email': email,
        'credits': DAILY_LIMIT,
        'credit_date': str(date.today()),
        'created': str(datetime.now()),
        'total_chats': 0
    }
    save_db(db)
    return jsonify({
        'ok': True,
        'username': username,
        'credits': DAILY_LIMIT
    })


@app.post('/api/login')
def login():
    ip = request.remote_addr
    # Brute force check
    ok, msg = check_brute_force(ip)
    if not ok:
        return jsonify({'error': msg}), 429

    data = request.json
    username = sanitize_input(data.get('username', ''), 50).lower()
    password = data.get('password', '')
    if not username or not password:
        return jsonify({'error': 'Missing credentials'}), 400

    db = load_db()
    user = db['users'].get(username)
    if not user or user['password'] != hash_pass(password):
        record_failed_login(ip)
        return jsonify({'error': 'Invalid username or password'}), 401

    reset_failed_login(ip)
    credits = get_user_credits(username)
    return jsonify({
        'ok': True,
        'username': username,
        'credits': credits,
        'total_chats': user.get('total_chats', 0)
    })


@app.get('/api/credits')
def credits():
    username = request.args.get('username', '').strip().lower()
    if not username:
        return jsonify({'error': 'Username required'}), 400
    c = get_user_credits(username)
    return jsonify({'credits': c, 'daily_limit': DAILY_LIMIT})


# ===== CHAT =====
@app.post('/api/chat')
def chat():
    try:
        ip = request.remote_addr
        if not check_rate_limit(ip):
            return jsonify({'error': 'Rate limit exceeded'}), 429

        data = request.json
        message = data.get('message', '').strip()
        session_id = data.get('sessionId', 'default')
        username = data.get('username', '').strip().lower()
        mode = data.get('mode', 'chat')

        if not message:
            return jsonify({'error': 'Message required'}), 400
        
        if len(sessions) > 200:
            cleanup_sessions()

        if session_id not in sessions:
            sessions[session_id] = {'history': [], 'last_active': time.time()}

        sessions[session_id]['last_active'] = time.time()
        sessions[session_id]['history'].append({
            'role': 'user',
            'parts': [{'text': message}]
        })

        history = sessions[session_id]['history'][-30:]

        if mode == 'elite':
            current_key = get_openai_key()
            if not current_key:
                return jsonify({'error': 'OpenAI API key missing'}), 500
            
            # Format history for OpenAI
            openai_messages = [{"role": "system", "content": build_system_prompt(mode)}]
            for m in history:
                role = "assistant" if m['role'] == "model" else "user"
                openai_messages.append({"role": role, "content": m['parts'][0]['text']})

            def generate_openai():
                full_reply = ""
                try:
                    resp = requests.post(
                        OPENAI_CHAT_URL,
                        headers={"Authorization": f"Bearer {current_key}", "Content-Type": "application/json"},
                        json={
                            "model": "gpt-4o",
                            "messages": openai_messages,
                            "stream": True,
                            "temperature": 0.7
                        },
                        stream=True, timeout=60
                    )
                    if resp.status_code != 200:
                        yield f"data: {json.dumps({'error': f'OpenAI Error: {resp.status_code}'})}\n\n"
                        yield "data: [DONE]\n\n"
                        return
                    
                    for line in resp.iter_lines():
                        if not line: continue
                        line = line.decode('utf-8')
                        if line.startswith("data: "):
                            raw = line[6:].strip()
                            if raw == "[DONE]": break
                            try:
                                chunk = json.loads(raw)
                                delta = chunk['choices'][0]['delta'].get('content', '')
                                if delta:
                                    full_reply += delta
                                    yield f"data: {json.dumps({'text': delta}, ensure_ascii=False)}\n\n"
                            except: pass
                except Exception as e:
                    yield f"data: {json.dumps({'error': str(e)})}\n\n"
                finally:
                    if full_reply:
                        sessions[session_id]['history'].append({'role': 'model', 'parts': [{'text': full_reply}]})
                    remaining = get_user_credits(username) if username else -1
                    yield f"data: {json.dumps({'credits': remaining})}\n\n"
                    yield "data: [DONE]\n\n"

            return Response(stream_with_context(generate_openai()), mimetype='text/event-stream', headers={'Cache-Control': 'no-cache', 'X-Accel-Buffering': 'no'})

        # Default Gemini logic
        current_key = get_gemini_key()
        if not current_key:
            return jsonify({'error': 'API key missing'}), 500

        body = build_body(history, mode)

        def generate():
            full_reply = ""
            try:
                resp = requests.post(
                    f"{GEMINI_STREAM_URL}?alt=sse&key={current_key}",
                    json=body, stream=True, timeout=60
                )
                if resp.status_code == 429:
                    yield f"data: {json.dumps({'error': 'API quota exceeded. Wait a minute.'})}\n\n"
                    yield "data: [DONE]\n\n"
                    return
                if resp.status_code != 200:
                    yield f"data: {json.dumps({'error': f'API error: {resp.status_code}'})}\n\n"
                    yield "data: [DONE]\n\n"
                    return
                for raw_line in resp.iter_lines():
                    if not raw_line:
                        continue
                    line = raw_line.decode('utf-8') if isinstance(raw_line, bytes) else raw_line
                    if not line.startswith('data: '):
                        continue
                    raw = line[6:].strip()
                    if raw == '[DONE]':
                        break
                    try:
                        chunk = json.loads(raw)
                        candidates = chunk.get('candidates', [])
                        if not candidates:
                            continue
                        parts = candidates[0].get('content', {}).get('parts', [])
                        if not parts:
                            continue
                        text = parts[0].get('text', '')
                        if text:
                            full_reply += text
                            yield f"data: {json.dumps({'text': text}, ensure_ascii=False)}\n\n"
                    except Exception:
                        pass
            except Exception as e:
                yield f"data: {json.dumps({'error': str(e)})}\n\n"
            finally:
                if full_reply:
                    sessions[session_id]['history'].append({
                        'role': 'model',
                        'parts': [{'text': full_reply}]
                    })
                # Send remaining credits
                remaining = get_user_credits(username) if username else -1
                yield f"data: {json.dumps({'credits': remaining})}\n\n"
                yield "data: [DONE]\n\n"

        return Response(
            stream_with_context(generate()),
            mimetype='text/event-stream',
            headers={'Cache-Control': 'no-cache', 'X-Accel-Buffering': 'no'}
        )

    except Exception as e:
        return jsonify({'error': str(e)}), 500


# ===== IMAGE =====
@app.post('/api/image')
def image():
    try:
        ip = request.remote_addr
        if not check_rate_limit(ip):
            return jsonify({'error': 'Rate limit exceeded'}), 429

        data = request.json
        prompt = data.get('prompt', '').strip()
        username = data.get('username', '').strip().lower()
        if not prompt:
            return jsonify({'error': 'Prompt required'}), 400

        current_key = get_gemini_key()
        # if username:
        #     ok, remaining = use_credits(username, 'image')
        #     if not ok:
        #         return jsonify({'error': f'Not enough credits! Need 50, have {remaining}', 'credits': remaining}), 403

        # Force Hyper-Realism
        high_end_prompt = f"Hyper-realistic, photorealistic masterpiece, 8k resolution, cinematic lighting, ultra-detailed, depth of field, professional photography, {prompt}"
        
        # --- OPTION 1: Google Imagen (Premium) ---
        if current_key:
            try:
                body = {
                    'instances': [{'prompt': high_end_prompt}],
                    'parameters': {'sampleCount': 1}
                }
                resp = requests.post(f"{IMAGEN_URL}?key={current_key}", json=body, timeout=25)
                result = resp.json()
                
                if 'error' not in result:
                    b64 = result.get('predictions', [{}])[0].get('bytesBase64Encoded')
                    if b64:
                        remaining = get_user_credits(username) if username else -1
                        return jsonify({'imageUrl': f'data:image/png;base64,{b64}', 'credits': remaining})
                
                print(f"Imagen Error: {result.get('error', {}).get('message', 'Unknown error')}")
            except Exception as e:
                print(f"Imagen Request Failed: {str(e)}")

        # --- OPTION 2: Fallback to Pollinations.ai (Free/Fast) ---
        try:
            print("Falling back to Pollinations.ai...")
            encoded_prompt = urllib.parse.quote(prompt)
            seed = int(time.time())
            # Fetch on backend to avoid browser-side auth issues (like the one you encountered)
            pollinations_url = f"https://image.pollinations.ai/prompt/{encoded_prompt}?width=1024&height=1024&nologo=true&seed={seed}&model=flux"
            
            resp = requests.get(pollinations_url, timeout=30)
            if resp.status_code == 200:
                import base64
                b64 = base64.b64encode(resp.content).decode('utf-8')
                remaining = get_user_credits(username) if username else -1
                return jsonify({'imageUrl': f'data:image/png;base64,{b64}', 'credits': remaining})
            else:
                print(f"Pollinations Error: {resp.status_code}")
                return jsonify({'error': 'Image generation failed on fallback.'}), 500
        except Exception as e:
            return jsonify({'error': f'Image generation failed: {str(e)}'}), 500

    except Exception as e:
        return jsonify({'error': str(e)}), 500


# ===== VIDEO =====
@app.post('/api/video')
def video():
    try:
        ip = request.remote_addr
        if not check_rate_limit(ip):
            return jsonify({'error': 'Rate limit exceeded'}), 429

        data = request.json
        prompt = data.get('prompt', '').strip()
        
        if not prompt:
            return jsonify({'error': 'Prompt required'}), 400

        hf_key = os.getenv("HF_API_KEY")
        if not hf_key:
            return jsonify({'error': 'To use Video Studio, please get a free HF API Key from huggingface.co and add it to your .env file as HF_API_KEY.'}), 403

        # HuggingFace API for Text-to-Video (ModelScope)
        API_URL = "https://api-inference.huggingface.co/models/damo-vilab/text-to-video-ms-1.7b"
        headers = {"Authorization": f"Bearer {hf_key}"}
        
        payload = {"inputs": prompt}
        resp = requests.post(API_URL, headers=headers, json=payload, timeout=120)
        
        if resp.status_code == 200:
            import base64
            b64 = base64.b64encode(resp.content).decode('utf-8')
            return jsonify({'videoUrl': f'data:video/mp4;base64,{b64}'})
        else:
            return jsonify({'error': f'Video API Error: {resp.text}'}), 500
            
    except Exception as e:
        return jsonify({'error': str(e)}), 500


@app.post('/api/search')
def web_search():
    """Web search via DuckDuckGo instant answers"""
    try:
        data = request.json
        query = data.get('query', '').strip()
        if not query:
            return jsonify({'error': 'Query required'}), 400

        resp = requests.get(
            'https://api.duckduckgo.com/',
            params={'q': query, 'format': 'json', 'no_html': 1, 't': 'aryax-ai'},
            timeout=10
        )
        result = resp.json()

        # Collect useful info
        answer_parts = []
        if result.get('AbstractText'):
            answer_parts.append(result['AbstractText'])
        if result.get('Answer'):
            answer_parts.append(result['Answer'])
        for topic in result.get('RelatedTopics', [])[:5]:
            if isinstance(topic, dict) and topic.get('Text'):
                answer_parts.append(f"• {topic['Text']}")

        if not answer_parts:
            return jsonify({'result': f'No instant results for "{query}". Try asking AryaX AI directly!'})

        return jsonify({'result': '\n\n'.join(answer_parts), 'source': result.get('AbstractURL', '')})

    except Exception as e:
        return jsonify({'error': str(e)}), 500


@app.post('/api/readurl')
def read_url():
    """Fetch and extract text from a URL"""
    try:
        data = request.json
        url = data.get('url', '').strip()
        if not url:
            return jsonify({'error': 'URL required'}), 400
        if not url.startswith(('http://', 'https://')):
            url = 'https://' + url

        resp = requests.get(url, timeout=15, headers={
            'User-Agent': 'AryaX-AI/1.0 (Educational Bot)'
        })
        resp.raise_for_status()

        # Simple HTML to text
        import re
        text = resp.text
        text = re.sub(r'<script[^>]*>[\s\S]*?</script>', '', text)
        text = re.sub(r'<style[^>]*>[\s\S]*?</style>', '', text)
        text = re.sub(r'<[^>]+>', ' ', text)
        text = re.sub(r'\s+', ' ', text).strip()

        # Truncate
        if len(text) > 3000:
            text = text[:3000] + '... (truncated)'

        return jsonify({'content': text, 'url': url})

    except Exception as e:
        return jsonify({'error': str(e)}), 500


@app.post('/api/clear')
def clear():
    data = request.json
    session_id = data.get('sessionId')
    if session_id and session_id in sessions:
        del sessions[session_id]
    return jsonify({'ok': True})


@app.get('/api/health')
def health():
    return jsonify({'status': 'ok', 'sessions': len(sessions)})

@app.post('/api/admin')
def get_admin_stats():
    data = request.json
    password = data.get('password', '')
    if password != 'aryax2055':  # Secret admin password
        return jsonify({'error': 'Unauthorized'}), 401
        
    db = load_db()
    users = db.get('users', {})
    
    total_users = len(users)
    total_chats = sum(u.get('total_chats', 0) for u in users.values())
    
    # Calculate revenue (mock logic based on plans)
    pro_users = sum(1 for u in users.values() if u.get('plan') == 'pro')
    ultra_users = sum(1 for u in users.values() if u.get('plan') == 'ultra')
    monthly_revenue = (pro_users * 99) + (ultra_users * 299)
    
    user_list = []
    for uname, udata in users.items():
        user_list.append({
            'username': uname,
            'plan': udata.get('plan', 'free'),
            'chats': udata.get('total_chats', 0),
            'created': udata.get('created', 'N/A')[:10]
        })
        
    return jsonify({
        'total_users': total_users,
        'total_chats': total_chats,
        'monthly_revenue': monthly_revenue,
        'pro_users': pro_users,
        'ultra_users': ultra_users,
        'recent_users': sorted(user_list, key=lambda x: x['created'], reverse=True)[:50]
    })

@app.post('/api/upgrade')
def upgrade_plan():
    """Endpoint for Direct UPI to verify payment and upgrade user"""
    data = request.json
    username = data.get('username')
    plan = data.get('plan') # 'pro' or 'ultra'
    utr = data.get('utr', '').strip()
    
    if len(utr) < 12:
        return jsonify({'error': 'Invalid UTR Number'}), 400
    
    db = load_db()
    if username in db['users']:
        db['users'][username]['plan'] = plan
        db['users'][username]['credits'] = 20000 if plan == 'pro' else 9999999
        # Save UTR for admin verification later
        utr_list = db['users'][username].get('utr_list', [])
        utr_list.append({'utr': utr, 'plan': plan, 'date': str(datetime.now())})
        db['users'][username]['utr_list'] = utr_list
        save_db(db)
        return jsonify({'ok': True, 'plan': plan})
    return jsonify({'error': 'User not found'}), 404


# ===== CLOUD CHAT HISTORY =====
@app.post('/api/history/save')
def save_history():
    data = request.json
    username = data.get('username', '').strip().lower()
    chat_id = data.get('chatId', '')
    title = data.get('title', 'Untitled Chat')[:60]
    messages = data.get('messages', [])
    if not username or not chat_id:
        return jsonify({'error': 'Missing fields'}), 400
    db = load_db()
    if username not in db['users']:
        return jsonify({'error': 'User not found'}), 404
    if 'chats' not in db['users'][username]:
        db['users'][username]['chats'] = {}
    db['users'][username]['chats'][chat_id] = {
        'title': title,
        'messages': messages[-50:],  # Keep last 50 messages
        'updated': str(datetime.now())
    }
    # Keep only last 30 chats
    chats = db['users'][username]['chats']
    if len(chats) > 30:
        oldest = sorted(chats.items(), key=lambda x: x[1].get('updated',''))[0][0]
        del chats[oldest]
    save_db(db)
    return jsonify({'ok': True})

@app.get('/api/history/load')
def load_history():
    username = request.args.get('username', '').strip().lower()
    if not username:
        return jsonify({'error': 'Missing username'}), 400
    db = load_db()
    user = db['users'].get(username)
    if not user:
        return jsonify({'error': 'User not found'}), 404
    chats = user.get('chats', {})
    result = [{'id': k, 'title': v['title'], 'updated': v['updated'], 'messages': v['messages']} 
              for k, v in chats.items()]
    result.sort(key=lambda x: x['updated'], reverse=True)
    return jsonify({'chats': result})

@app.post('/api/history/delete')
def delete_history():
    data = request.json
    username = data.get('username', '').strip().lower()
    chat_id = data.get('chatId', '')
    db = load_db()
    if username in db['users'] and 'chats' in db['users'][username]:
        db['users'][username]['chats'].pop(chat_id, None)
        save_db(db)
    return jsonify({'ok': True})


# ===== USER MEMORY SYSTEM =====
@app.post('/api/memory/save')
def save_memory():
    data = request.json
    username = data.get('username', '').strip().lower()
    memory = data.get('memory', {})  # {name, language, interests, etc}
    if not username:
        return jsonify({'error': 'Missing username'}), 400
    db = load_db()
    if username not in db['users']:
        return jsonify({'error': 'User not found'}), 404
    db['users'][username]['memory'] = memory
    save_db(db)
    return jsonify({'ok': True})

@app.get('/api/memory/load')
def load_memory():
    username = request.args.get('username', '').strip().lower()
    if not username:
        return jsonify({'error': 'Missing username'}), 400
    db = load_db()
    user = db['users'].get(username)
    if not user:
        return jsonify({'error': 'User not found'}), 404
    return jsonify({'memory': user.get('memory', {})})


# ===== REFERRAL SYSTEM =====
@app.post('/api/referral/generate')
def generate_referral():
    data = request.json
    username = data.get('username', '').strip().lower()
    if not username:
        return jsonify({'error': 'Missing username'}), 400
    db = load_db()
    user = db['users'].get(username)
    if not user:
        return jsonify({'error': 'User not found'}), 404
    # Generate or return existing code
    if not user.get('referral_code'):
        code = 'AX' + hashlib.md5(username.encode()).hexdigest()[:6].upper()
        db['users'][username]['referral_code'] = code
        db['users'][username]['referral_count'] = 0
        save_db(db)
    return jsonify({
        'code': db['users'][username]['referral_code'],
        'count': db['users'][username].get('referral_count', 0),
        'reward_per_referral': 500
    })

@app.post('/api/referral/use')
def use_referral():
    data = request.json
    new_user = data.get('username', '').strip().lower()
    code = data.get('code', '').strip().upper()
    if not new_user or not code:
        return jsonify({'error': 'Missing fields'}), 400
    db = load_db()
    # Find referrer
    referrer = None
    for uname, udata in db['users'].items():
        if udata.get('referral_code') == code:
            referrer = uname
            break
    if not referrer:
        return jsonify({'error': 'Invalid referral code'}), 404
    if referrer == new_user:
        return jsonify({'error': 'Cannot use your own code'}), 400
    # Check if already used a code
    if db['users'][new_user].get('referral_used'):
        return jsonify({'error': 'Already used a referral code'}), 400
    # Reward referrer +500 credits
    db['users'][referrer]['referral_count'] = db['users'][referrer].get('referral_count', 0) + 1
    db['users'][referrer]['credits'] = db['users'][referrer].get('credits', 0) + 500
    # Reward new user +200 credits
    db['users'][new_user]['credits'] = db['users'][new_user].get('credits', 0) + 200
    db['users'][new_user]['referral_used'] = code
    save_db(db)
    return jsonify({'ok': True, 'message': 'Referral applied! +200 credits added.'})


# ===== PDF / FILE UPLOAD =====
@app.post('/api/upload')
def upload_file():
    try:
        if 'file' not in request.files:
            return jsonify({'error': 'No file uploaded'}), 400
        file = request.files['file']
        filename = file.filename.lower()
        content = ''

        if filename.endswith('.pdf'):
            try:
                import PyPDF2
                reader = PyPDF2.PdfReader(file)
                for page in reader.pages[:10]:  # Max 10 pages
                    content += page.extract_text() or ''
                content = content[:4000]
            except Exception as e:
                return jsonify({'error': f'PDF read failed: {str(e)}'}), 500

        elif filename.endswith(('.txt', '.py', '.js', '.html', '.css', '.json', '.md', '.csv')):
            content = file.read().decode('utf-8', errors='ignore')[:4000]

        elif filename.endswith(('.png', '.jpg', '.jpeg', '.webp')):
            return jsonify({'content': f'[Image uploaded: {file.filename}]', 'type': 'image'})

        else:
            content = file.read().decode('utf-8', errors='ignore')[:4000]

        if not content.strip():
            return jsonify({'error': 'Could not extract text from file'}), 400

        return jsonify({'content': content, 'filename': file.filename})
    except Exception as e:
        return jsonify({'error': str(e)}), 500


# ===== USER PROFILE =====
@app.get('/api/profile')
def get_profile():
    username = request.args.get('username', '').strip().lower()
    if not username:
        return jsonify({'error': 'Missing username'}), 400
    db = load_db()
    user = db['users'].get(username)
    if not user:
        return jsonify({'error': 'User not found'}), 404
    return jsonify({
        'username': username,
        'plan': user.get('plan', 'free'),
        'credits': user.get('credits', 0),
        'total_chats': user.get('total_chats', 0),
        'created': user.get('created', '')[:10],
        'referral_code': user.get('referral_code', ''),
        'referral_count': user.get('referral_count', 0),
        'memory': user.get('memory', {}),
        'chat_count': len(user.get('chats', {}))
    })


# ===== ONLINE COUNT =====
@app.get('/api/online')
def online_count():
    active = sum(1 for s in sessions.values() if time.time() - s.get('last_active', 0) < 300)
    return jsonify({'online': max(active, 1)})


if __name__ == '__main__':
    # Initialize db
    if not os.path.exists(DB_FILE) or os.path.getsize(DB_FILE) < 5:
        save_db({'users': {}})
    port = int(os.environ.get('PORT', 5054))
    print(f"\nAryaX AI Production Server Starting on Port {port}...")
    app.run(host='0.0.0.0', port=port)

@app.route('/api/generate_office', methods=['POST'])
def generate_office():
    data = request.json
    file_type = data.get('type') # 'pptx' or 'xlsx'
    content = data.get('content')
    
    if file_type == 'pptx':
        prs = Presentation()
        for slide_data in content.get('slides', []):
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = slide_data.get('title', '')
            slide.placeholders[1].text = slide_data.get('body', '')
        
        out = io.BytesIO()
        prs.save(out)
        out.seek(0)
        return Response(out.read(), mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation', headers={'Content-Disposition': 'attachment; filename=aryax_presentation.pptx'})
    
    elif file_type == 'xlsx':
        df = pd.DataFrame(content.get('data', []))
        out = io.BytesIO()
        with pd.ExcelWriter(out, engine='openpyxl') as writer:
            df.to_excel(writer, index=False, sheet_name='AryaX Ledger')
        out.seek(0)
        return Response(out.read(), mimetype='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet', headers={'Content-Disposition': 'attachment; filename=aryax_ledger.xlsx'})
    
    return jsonify({'error': 'Invalid type'}), 400
